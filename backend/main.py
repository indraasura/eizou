from fastapi import FastAPI, UploadFile, File, Form, HTTPException, Depends, Header
from fastapi.responses import HTMLResponse
from fastapi.middleware.cors import CORSMiddleware
import os
import pandas as pd
from pptx import Presentation
from dotenv import load_dotenv, find_dotenv
import io
from supabase import create_client, Client
import PyPDF2
import boto3
import re
import json

# --- LANGCHAIN IMPORTS ---
from langchain_text_splitters import RecursiveCharacterTextSplitter
from langchain_google_genai import GoogleGenerativeAIEmbeddings, ChatGoogleGenerativeAI
from langchain_aws import ChatBedrockConverse, BedrockEmbeddings
from langchain_core.prompts import ChatPromptTemplate
from langchain_core.documents import Document
from langchain_core.output_parsers import StrOutputParser

load_dotenv(find_dotenv())

app = FastAPI()

app.add_middleware(
    CORSMiddleware,
    allow_origins=["*"],
    allow_credentials=False,
    allow_methods=["*"],
    allow_headers=["*"],
)

# --- SUPABASE INIT ---
url: str = os.getenv("SUPABASE_URL")
key: str = os.getenv("SUPABASE_KEY")
supabase: Client = create_client(url, key)

# --- AUTH & RBAC DEPENDENCIES ---
def get_current_user(authorization: str = Header(None)):
    if not authorization or not authorization.startswith("Bearer "):
        raise HTTPException(status_code=401, detail="Missing Token")
    
    token = authorization.split(" ")[1]
    
    try:
        temp_client = create_client(os.getenv("SUPABASE_URL"), os.getenv("SUPABASE_KEY"))
        user_response = temp_client.auth.get_user(token)
        user_id = user_response.user.id
        email = user_response.user.email
        
        profile = supabase.table("profiles").select("role").eq("id", user_id).execute()
        role = profile.data[0]["role"] if profile.data else "user"
        
        return {"id": user_id, "email": email, "role": role}
    except Exception as e:
        print(f"Auth verification error: {str(e)}")
        raise HTTPException(status_code=401, detail="Invalid or Expired Token")

def require_admin(user: dict = Depends(get_current_user)):
    if user.get("role") != "admin":
        raise HTTPException(status_code=403, detail="Admin access required")
    return user

# --- AUTH ENDPOINTS ---
@app.post("/login")
def login(email: str = Form(...), password: str = Form(...)):
    try:
        temp_client = create_client(os.getenv("SUPABASE_URL"), os.getenv("SUPABASE_KEY"))
        res = temp_client.auth.sign_in_with_password({"email": email, "password": password})
        token = res.session.access_token
        
        profile = supabase.table("profiles").select("role").eq("id", res.user.id).execute()
        role = profile.data[0]["role"] if profile.data else "user"
        
        return {"token": token, "role": role, "email": res.user.email}
    except Exception:
        raise HTTPException(status_code=401, detail="Invalid credentials")

# --- ADMIN ENDPOINTS ---
@app.post("/admin/users")
def create_user(email: str = Form(...), password: str = Form(...), role: str = Form(...), admin: dict = Depends(require_admin)):
    try:
        admin_client = create_client(os.getenv("SUPABASE_URL"), os.getenv("SUPABASE_KEY"))
        
        res = admin_client.auth.admin.create_user({
            "email": email,
            "password": password,
            "email_confirm": True
        })
        new_user_id = res.user.id
        
        admin_client.table("profiles").update({"role": role}).eq("id", new_user_id).execute()
        return {"status": "User provisioned successfully"}
    except Exception as e:
        raise HTTPException(status_code=400, detail=str(e))

@app.post("/admin/projects")
def create_project(name: str = Form(...), admin: dict = Depends(require_admin)):
    res = supabase.table("projects").insert({"name": name}).execute()
    return {"status": "Project created", "data": res.data}

@app.post("/admin/assign")
def assign_user(user_email: str = Form(...), project_id: int = Form(...), admin: dict = Depends(require_admin)):
    profile = supabase.table("profiles").select("id").eq("email", user_email).execute()
    if not profile.data:
        raise HTTPException(status_code=404, detail="User not found")
    
    user_id = profile.data[0]["id"]
    try:
        supabase.table("project_users").insert({"project_id": project_id, "user_id": user_id}).execute()
        return {"status": "User assigned"}
    except Exception:
        return {"status": "User already assigned"}

@app.get("/admin/list_users")
def list_users(admin: dict = Depends(require_admin)):
    profiles = supabase.table("profiles").select("*").execute()
    return {"users": profiles.data}

@app.get("/admin/files")
def list_files(admin: dict = Depends(require_admin)):
    try:
        files = supabase.table("project_files").select("*").execute()
        return {"files": files.data}
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))

@app.delete("/admin/projects/{project_id}")
def delete_project(project_id: int, admin: dict = Depends(require_admin)):
    try:
        supabase.table("projects").delete().eq("id", project_id).execute()
        return {"status": "Project deleted"}
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))

@app.put("/admin/projects/{project_id}")
def edit_project(project_id: int, name: str = Form(...), admin: dict = Depends(require_admin)):
    try:
        supabase.table("projects").update({"name": name}).eq("id", project_id).execute()
        return {"status": "Project updated"}
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))

@app.delete("/admin/users/{user_id}")
def delete_user(user_id: str, admin: dict = Depends(require_admin)):
    try:
        supabase.auth.admin.delete_user(user_id)
        supabase.table("profiles").delete().eq("id", user_id).execute()
        return {"status": "User deleted"}
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))

@app.put("/admin/users/{user_id}/role")
def update_user_role(user_id: str, role: str = Form(...), admin: dict = Depends(require_admin)):
    try:
        supabase.table("profiles").update({"role": role}).eq("id", user_id).execute()
        return {"status": "Role updated"}
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))

@app.delete("/admin/files/{file_id}")
def delete_file(file_id: int, admin: dict = Depends(require_admin)):
    try:
        file_record = supabase.table("project_files").select("*").eq("id", file_id).execute()
        if file_record.data:
            path = f"project_{file_record.data[0]['project_id']}/{file_record.data[0]['file_name']}"
            supabase.storage.from_("project_files").remove([path])
        
        supabase.table("project_files").delete().eq("id", file_id).execute()
        return {"status": "File deleted"}
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))

# --- USER ENDPOINTS ---
@app.get("/projects")
def get_user_projects(user: dict = Depends(get_current_user)):
    try:
        if user["role"] == "admin":
            projects = supabase.table("projects").select("*").execute()
            return {"projects": projects.data}
        else:
            assignments = supabase.table("project_users").select("project_id").eq("user_id", user["id"]).execute()
            
            if not assignments.data:
                return {"projects": []}
                
            project_ids = [str(a["project_id"]) for a in assignments.data]
            projects = supabase.table("projects").select("*").in_("id", project_ids).execute()
            
            return {"projects": projects.data}
            
    except Exception as e:
        print(f"❌ Error fetching projects: {str(e)}")
        raise HTTPException(status_code=500, detail=f"Database error: {str(e)}")

# --- AI DATA ENDPOINTS ---
@app.post("/upload")
async def upload_files(files: list[UploadFile] = File(...), project_id: int = Form(...), model: str = Form(...), admin: dict = Depends(require_admin)):
    documents = []
    
    for file in files:
        contents = await file.read()
        ext = file.filename.split('.')[-1].lower()
        file_path = f"project_{project_id}/{file.filename}"
        
        file_url = "#" 
        
        try:
            supabase.storage.from_("project_files").upload(file_path, contents, file_options={"upsert": "true"})
            file_url = supabase.storage.from_("project_files").get_public_url(file_path)
            
            supabase.table("project_files").insert({
                "project_id": project_id,
                "file_name": file.filename,
                "file_url": file_url
            }).execute()
        except Exception as e:
            print(f"⚠️ Storage Error for {file.filename}: {str(e)}")

        raw_text = ""
        try:
            if ext in ['xlsx', 'xls']:
                raw_text = pd.read_excel(io.BytesIO(contents)).to_markdown()
            elif ext == 'csv':
                raw_text = pd.read_csv(io.BytesIO(contents)).to_markdown()
            elif ext in ['pptx', 'ppt']:
                prs = Presentation(io.BytesIO(contents))
                raw_text = "\n".join([shape.text for slide in prs.slides for shape in slide.shapes if hasattr(shape, "text")])
            elif ext == 'pdf':
                pdf_reader = PyPDF2.PdfReader(io.BytesIO(contents))
                raw_text = "\n".join([page.extract_text() for page in pdf_reader.pages if page.extract_text()])
            else:
                raw_text = contents.decode('utf-8', errors='ignore')
        except Exception as e:
            print(f"⚠️ Parsing Error for {file.filename}: {str(e)}")

        if raw_text.strip():
            documents.append(Document(
                page_content=raw_text, 
                metadata={
                    "source": file.filename, 
                    "project_id": project_id, 
                    "file_url": file_url
                }
            ))
            
    if not documents:
        raise HTTPException(status_code=400, detail="No readable text found in the uploaded files.")

    text_splitter = RecursiveCharacterTextSplitter(chunk_size=2000, chunk_overlap=200)
    chunks = text_splitter.split_documents(documents)
    
    if "gemini" in model.lower():
        embeddings = GoogleGenerativeAIEmbeddings(
            model="gemini-embedding-001", 
            google_api_key=os.getenv("GOOGLE_API_KEY"), 
            transport="rest"
        )
    else:
        bedrock_client = boto3.client(
            service_name="bedrock-runtime",
            region_name=os.getenv("AWS_DEFAULT_REGION", "us-east-1")
        )
        embeddings = BedrockEmbeddings(client=bedrock_client, model_id="amazon.titan-embed-text-v2:0")
        
    try:
        texts = [chunk.page_content for chunk in chunks]
        metadatas = [chunk.metadata for chunk in chunks]
        
        raw_vectors = embeddings.embed_documents(texts)
        truncated_vectors = [vec[:1024] for vec in raw_vectors]
        
        records = []
        for text, meta, vec in zip(texts, metadatas, truncated_vectors):
            records.append({
                "content": text,
                "metadata": meta,
                "embedding": vec
            })
            
        supabase.table("project_documents").insert(records).execute()
        
    except Exception as e:
        print(f"❌ Vector DB Error: {str(e)}")
        raise HTTPException(status_code=500, detail=f"Failed to sync with AI memory: {str(e)}")
    
    return {"status": "success", "message": f"Knowledge base synced. {len(files)} files processed."}

@app.post("/chat")
async def chat(message: str = Form(...), project_id: int = Form(...), model: str = Form(...), user: dict = Depends(get_current_user)):
    
    if "gemini" in model.lower():
        embeddings = GoogleGenerativeAIEmbeddings(model="gemini-embedding-001", google_api_key=os.getenv("GOOGLE_API_KEY"), transport="rest")
        llm = ChatGoogleGenerativeAI(model=model, google_api_key=os.getenv("GOOGLE_API_KEY"), temperature=0, transport="rest")
    else:
        bedrock_client = boto3.client(
            service_name="bedrock-runtime",
            region_name=os.getenv("AWS_DEFAULT_REGION", "ap-south-1"),
            aws_access_key_id=os.getenv("AWS_ACCESS_KEY_ID"),
            aws_secret_access_key=os.getenv("AWS_SECRET_ACCESS_KEY")
        )
        embeddings = BedrockEmbeddings(client=bedrock_client, model_id="amazon.titan-embed-text-v2:0")
        llm = ChatBedrockConverse(client=bedrock_client, model_id="openai.gpt-oss-20b-1:0", temperature=0)
    
    try:
        query_embedding = embeddings.embed_query(message)[:1024]
        rpc_response = supabase.rpc("match_project_documents", {
            "query_embedding": query_embedding,
            "match_count": 20,
            "filter": {"project_id": int(project_id)} 
        }).execute()
        
        chunks = rpc_response.data or []
        context_text = ""
        potential_filenames = set()

        for row in chunks:
            meta = row.get("metadata", {})
            name = meta.get("source", "Unknown")
            content = row.get("content", "")
            context_text += f"\n---\nDOCUMENT: {name}\nCONTENT: {content}\n"
            potential_filenames.add(name)

    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Search error: {str(e)}")
    
    system_prompt = (
        "You are an enterprise data assistant. Answer based strictly on the context provided.\n\n"
        "### OUTPUT RULES:\n"
        "1. VISUALS: If the user asks for a graph, trend, or comparison, you MUST generate a valid JSON object "
        "wrapped in a ```chart ... ``` block. Use 'bar', 'line', or 'pie'.\n"
        "Structure: {{ \"type\": \"bar\", \"data\": {{ \"labels\": [\"Jan\", \"Feb\"], \"datasets\": [{{ \"label\": \"Revenue\", \"data\": [100, 200], \"backgroundColor\": \"#0A56D0\" }}] }}, \"options\": {{ \"responsive\": true }} }}\n\n"
        "2. CITATIONS: At the very end of your response, after any text or charts, you MUST list the document names used.\n"
        "Format: SOURCES: [File1.pdf, File2.xlsx]\n\n"
        "Context: {context}"
    )

    prompt_template = ChatPromptTemplate.from_messages([
        ("system", system_prompt), ("human", "{input}"),
    ])
    
    rag_chain = prompt_template | llm | StrOutputParser()
    
    try:
        full_response = rag_chain.invoke({"context": context_text, "input": message})
        
        cited_sources = []
        clean_answer = full_response
        
        # --- 🚨 FULL-PROOF SOURCE EXTRACTION ENGINE 🚨 ---
        # Hunts for variations like "SOURCES:", "Sources:", "**SOURCES:**", etc.
        source_match = re.search(r'\*?\*[Ss]ources?:?\*?\*?\s*(.*)', full_response, re.IGNORECASE | re.DOTALL)
        
        if source_match:
            # Cleanly severe the sources section from the visible chat answer
            clean_answer = full_response[:source_match.start()].strip()
            
            raw_sources_str = source_match.group(1).strip()
            # Strip out any brackets or parentheses the AI might mistakenly add
            clean_sources_str = re.sub(r'[\[\]\(\)]', '', raw_sources_str)
            
            if clean_sources_str:
                # Split by comma and strip errant whitespace/quotes
                names_list = [n.strip().strip('"\'') for n in clean_sources_str.split(",") if n.strip()]
                
                if names_list:
                    # Fetch real URL mappings from Supabase
                    file_data = supabase.table("project_files") \
                        .select("file_name, file_url") \
                        .in_("file_name", names_list) \
                        .eq("project_id", project_id) \
                        .execute()
                    
                    url_map = {f["file_name"]: f["file_url"] for f in file_data.data}
                    
                    for name in names_list:
                        cited_sources.append({
                            "name": name, 
                            "url": url_map.get(name, "#")
                        })

        return {"answer": clean_answer, "sources": cited_sources}
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))
    

# Gitbook iframe embed
@app.get("/embed", response_class=HTMLResponse)
async def embed_widget(project_id: int = 999):
    """
    Serves a standalone, full-frame chat UI for embedding in GitBook or other platforms.
    Usage: https://your-domain.com/embed?project_id=999
    """
    html_template = """
    <!DOCTYPE html>
    <html lang="en">
    <head>
        <meta charset="UTF-8">
        <meta name="viewport" content="width=device-width, initial-scale=1.0">
        <title>Ask Nexus</title>
        <style>
            body { 
                margin: 0; 
                font-family: -apple-system, BlinkMacSystemFont, "Segoe UI", Roboto, Helvetica, Arial, sans-serif; 
                display: flex; flex-direction: column; height: 100vh; 
                background: #FFFFFF; overflow: hidden;
            }
            #header { 
                background: #0A56D0; color: white; padding: 14px 20px; 
                font-weight: 600; font-size: 15px; display: flex; align-items: center; gap: 10px;
                box-shadow: 0 2px 10px rgba(0,0,0,0.1); z-index: 10;
            }
            #chat { 
                flex: 1; overflow-y: auto; padding: 20px; background: #F8F9FA; 
                display: flex; flex-direction: column; gap: 16px; scroll-behavior: smooth;
            }
            .msg { 
                padding: 12px 16px; border-radius: 16px; max-width: 90%; 
                font-size: 14.5px; line-height: 1.5; word-wrap: break-word; 
            }
            .msg.bot { 
                background: white; border: 1px solid #E3E3E3; align-self: flex-start; 
                color: #1F1F1F; border-bottom-left-radius: 4px; box-shadow: 0 2px 5px rgba(0,0,0,0.02);
            }
            .msg.user { 
                background: #0A56D0; color: white; align-self: flex-end; 
                border-bottom-right-radius: 4px; box-shadow: 0 2px 5px rgba(10,86,208,0.2);
            }
            .sources { 
                margin-top: 12px; padding-top: 12px; border-top: 1px solid #E3E3E3; font-size: 12px; 
            }
            .source-link { 
                color: #0A56D0; text-decoration: none; font-weight: 500; display: block; margin-top: 6px; 
            }
            .source-link:hover { text-decoration: underline; }
            #input-area { 
                padding: 16px; background: white; border-top: 1px solid #E3E3E3; 
                display: flex; gap: 12px; align-items: center;
            }
            #input { 
                flex: 1; padding: 12px 16px; border: 1px solid #E3E3E3; border-radius: 24px; 
                outline: none; font-size: 14.5px; transition: border 0.2s; 
            }
            #input:focus { border-color: #0A56D0; }
            #send { 
                background: #0A56D0; color: white; border: none; border-radius: 50%; 
                width: 44px; height: 44px; cursor: pointer; display: flex; align-items: center; 
                justify-content: center; transition: background 0.2s; flex-shrink: 0;
            }
            #send:hover { background: #0842A0; }
            #send:disabled { background: #A8C7FA; cursor: not-allowed; }
            .loading { display: flex; gap: 4px; padding: 4px 8px; }
            .dot { width: 6px; height: 6px; background: #0A56D0; border-radius: 50%; animation: bounce 1.4s infinite ease-in-out both; }
            .dot:nth-child(1) { animation-delay: -0.32s; }
            .dot:nth-child(2) { animation-delay: -0.16s; }
            @keyframes bounce { 0%, 80%, 100% { transform: scale(0); } 40% { transform: scale(1); } }
        </style>
    </head>
    <body>
        <div id="header">
            <svg width="20" height="20" viewBox="0 0 24 24" fill="none" stroke="currentColor" stroke-width="2" stroke-linecap="round" stroke-linejoin="round"><path d="M21 15a2 2 0 0 1-2 2H7l-4 4V5a2 2 0 0 1 2-2h14a2 2 0 0 1 2 2z"></path></svg>
            Ask Nexus AI
        </div>
        <div id="chat">
            <div class="msg bot">Hello! I am trained on this documentation. What are you looking for?</div>
        </div>
        <div id="input-area">
            <input type="text" id="input" placeholder="Ask a question..." autocomplete="off" />
            <button id="send">
                <svg width="18" height="18" viewBox="0 0 24 24" fill="none" stroke="currentColor" stroke-width="2" stroke-linecap="round" stroke-linejoin="round"><line x1="22" y1="2" x2="11" y2="13"></line><polygon points="22 2 15 22 11 13 2 9 22 2"></polygon></svg>
            </button>
        </div>
        <script>
            const PROJECT_ID = __PROJECT_ID__;
            const chat = document.getElementById('chat');
            const input = document.getElementById('input');
            const sendBtn = document.getElementById('send');

            function appendMsg(text, sender, sources = []) {
                const msg = document.createElement('div');
                msg.className = 'msg ' + sender;
                
                let formattedText = text.replace(/\\*\\*(.*?)\\*\\*/g, '<strong>$1</strong>');
                
                let sourceHtml = '';
                if (sources && sources.length > 0) {
                    sourceHtml = '<div class="sources"><strong>Documents Referenced:</strong>';
                    sources.forEach(src => {
                        let url = (src.url && src.url !== "#") ? src.url : "javascript:void(0)";
                        let style = (url === "javascript:void(0)") ? "style='color:#666;text-decoration:none;cursor:default;'" : "";
                        sourceHtml += `<a href="${url}" target="_blank" class="source-link" ${style}>↳ ${src.name}</a>`;
                    });
                    sourceHtml += '</div>';
                }

                msg.innerHTML = formattedText + sourceHtml;
                chat.appendChild(msg);
                chat.scrollTop = chat.scrollHeight;
            }

            async function handleSend() {
                const text = input.value.trim();
                if (!text) return;

                appendMsg(text, 'user');
                input.value = '';
                sendBtn.disabled = true;

                const loader = document.createElement('div');
                loader.className = 'msg bot';
                loader.innerHTML = '<div class="loading"><div class="dot"></div><div class="dot"></div><div class="dot"></div></div>';
                chat.appendChild(loader);
                chat.scrollTop = chat.scrollHeight;

                const formData = new FormData();
                formData.append("message", text);
                formData.append("project_id", PROJECT_ID);
                formData.append("model", "gemini-2.5-flash");

                try {
                    const res = await fetch('/chat', { method: 'POST', body: formData });
                    const data = await res.json();
                    
                    chat.removeChild(loader);

                    if (res.ok) {
                        appendMsg(data.answer, 'bot', data.sources);
                    } else {
                        appendMsg('⚠️ Error: ' + (data.detail || 'Failed to analyze data.'), 'bot');
                    }
                } catch (e) {
                    chat.removeChild(loader);
                    appendMsg('❌ Network error. Check server connection.', 'bot');
                } finally {
                    sendBtn.disabled = false;
                    input.focus();
                }
            }

            sendBtn.addEventListener('click', handleSend);
            input.addEventListener('keypress', (e) => { if (e.key === 'Enter') handleSend(); });
        </script>
    </body>
    </html>
    """
    
    # Safely inject the project ID without f-string escaping nightmares
    final_html = html_template.replace("__PROJECT_ID__", str(project_id))
    return final_html

